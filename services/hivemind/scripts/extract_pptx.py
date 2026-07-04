#!/usr/bin/env python3
"""Extract plain text from a .pptx file. Reads the file path as argv[1], writes
extracted text to stdout. Used by hivemind (Rust) via subprocess, mirroring the
ffmpeg-subprocess pattern in transcription-service — python-pptx is the mature,
de facto standard for OOXML PowerPoint parsing; no Rust equivalent is close.
"""
import sys

from pptx import Presentation


def extract(path: str) -> str:
    prs = Presentation(path)
    lines = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {slide_index} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs)
                    if text.strip():
                        lines.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        lines.append("\t".join(cells))
            if shape.has_chart:
                chart = shape.chart
                title = ""
                try:
                    if chart.has_title:
                        title = chart.chart_title.text_frame.text
                except Exception:
                    pass
                if title:
                    lines.append(f"[Chart: {title}]")
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            lines.append(f"[Notes: {slide.notes_slide.notes_text_frame.text.strip()}]")
    return "\n".join(lines)


if __name__ == "__main__":
    if len(sys.argv) != 2:
        print("usage: extract_pptx.py <path-to-pptx>", file=sys.stderr)
        sys.exit(2)
    try:
        sys.stdout.write(extract(sys.argv[1]))
    except Exception as e:
        print(f"pptx extraction failed: {e}", file=sys.stderr)
        sys.exit(1)
